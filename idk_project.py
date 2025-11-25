
import os
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Image
from reportlab.lib.pagesizes import A4
from reportlab.lib.styles import getSampleStyleSheet
from reportlab.lib.units import inch
from docx import Document
from pptx import Presentation
from pptx.util import Inches
import matplotlib.pyplot as plt
import zipfile

# Prepare directories
os.makedirs('/mnt/data/project_package', exist_ok=True)

# ----- Create graphs -----
plt.figure()
plt.plot([1,2,3,4],[10,20,15,30])
plt.title("Sample Graph")
graph_path = "/mnt/data/project_package/graph.png"
plt.savefig(graph_path)
plt.close()

# ----- Create PDF -----
pdf_path = "/mnt/data/project_package/report.pdf"
styles = getSampleStyleSheet()
story = []
story.append(Paragraph("Machine Learning Based Intrusion Detection System (IDS) Report", styles['Title']))
story.append(Spacer(1, 12))
story.append(Paragraph("This is a generated sample report with graphs included.", styles['BodyText']))
story.append(Spacer(1, 12))
story.append(Image(graph_path, width=4*inch, height=3*inch))
doc = SimpleDocTemplate(pdf_path, pagesize=A4)
doc.build(story)

# ----- Create DOCX -----
docx_path = "/mnt/data/project_package/report.docx"
docx = Document()
docx.add_heading("Machine Learning Based IDS Report", level=1)
docx.add_paragraph("This is a generated sample report with graphs included.")
docx.add_picture(graph_path, width=Inches(4))
docx.save(docx_path)

# ----- Create PPTX -----
ppt_path = "/mnt/data/project_package/report.pptx"
prs = Presentation()
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
body = slide.placeholders[1]
title.text = "Machine Learning Based IDS"
body.text = "This presentation includes sample content and graphs."
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Graph"
slide.shapes.add_picture(graph_path, Inches(1), Inches(1), width=Inches(6))
prs.save(ppt_path)

# ----- Create ZIP -----
zip_path = "/mnt/data/IDS_Project_Full_Package.zip"
with zipfile.ZipFile(zip_path, 'w') as zipf:
    zipf.write(pdf_path, arcname="report.pdf")
    zipf.write(docx_path, arcname="report.docx")
    zipf.write(ppt_path, arcname="report.pptx")
    zipf.write(graph_path, arcname="graph.png")

zip_path

